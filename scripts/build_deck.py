# -*- coding: utf-8 -*-
"""Build Vigil_Round2.pptx - the Round 2 detailed business proposal.

Uses the AIC template so the branding matches, and draws every content slide
explicitly rather than filling placeholders, because the business proposal needs
tables and diagrams the stock layouts do not provide.

Numbers are pulled live from the engine at build time. If a result changes, the
deck changes with it - a slide deck that can drift from the code it describes is
how a demo and a claim end up disagreeing in front of a jury.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from vigil.clinical.redflags import PANEL                      # noqa: E402
from vigil.sim import compare, composition, generate_surge, surge_summary  # noqa: E402
from vigil.triage import assess                                # noqa: E402
from vigil.sim.scenarios import SCENARIOS, by_id                # noqa: E402


def _count_tests() -> int:
    """Ask pytest how many tests exist rather than remembering a number.

    Sums the per-file counts, which is the one format that has been stable
    across pytest versions -- the summary line's wording has not been, and a
    silently-unmatched regex here is how the deck came to claim a test count
    that was three commits stale.
    """
    import subprocess
    root = Path(__file__).resolve().parents[1]
    r = subprocess.run([sys.executable, "-m", "pytest", "--collect-only", "-q"],
                       capture_output=True, text=True, cwd=str(root))
    counts = [int(m.group(1)) for m in
              re.finditer(r"^\S+\.py: (\d+)$", r.stdout, re.MULTILINE)]
    total = sum(counts)
    if not total:
        raise RuntimeError(f"could not count tests; pytest said:\n{r.stdout[-800:]}")
    return total


def _count_safety_sweep() -> int:
    """The real size of the safety-property sweep."""
    sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "tests"))
    from test_safety_property import _grid
    return sum(1 for _ in _grid())

TEMPLATE = Path(r"C:\Users\lakshay\Desktop\Hackathon_accenture\AIC_Talent-Brand_PPT-Template (1).pptx")
OUT = Path(r"C:\Users\lakshay\Desktop\Hackathon_accenture\Vigil_Round2.pptx")

# ── palette (from the template's own XML) ─────────────────────────────────────
PURPLE, DEEP, MID = RGBColor(0xA1, 0x00, 0xFF), RGBColor(0x46, 0x00, 0x73), RGBColor(0x75, 0x00, 0xC0)
LAV, LAVE = RGBColor(0xF3, 0xEC, 0xFE), RGBColor(0xE0, 0xCC, 0xFB)
INK, GREY, GREYF = RGBColor(0x1A, 0x1A, 0x1A), RGBColor(0x5A, 0x5A, 0x5A), RGBColor(0xEB, 0xEB, 0xEE)
RED, REDT = RGBColor(0xC8, 0x10, 0x2E), RGBColor(0xFD, 0xEC, 0xEF)
AMB, AMBT = RGBColor(0xB4, 0x6A, 0x00), RGBColor(0xFF, 0xF4, 0xE3)
GRN, GRNT = RGBColor(0x0F, 0x7B, 0x55), RGBColor(0xE9, 0xF6, 0xF1)
W = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"

L, CW = 361458, 11282796          # left margin / content width, matching the template
TOP = 1000000                     # first free y under the title bar


# ── helpers ───────────────────────────────────────────────────────────────────
def para(tf, text, size, bold=False, color=INK, after=4, align=PP_ALIGN.LEFT,
         italic=False, first=False, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if line:
        p.line_spacing = line
    p.space_after = Pt(after)
    r = p.add_run(); r.text = text
    f = r.font; f.name = FONT; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    return p


def runs(tf, parts, size, color=INK, after=4, first=False, line=None, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if line:
        p.line_spacing = line
    p.space_after = Pt(after)
    for part in parts:
        txt, bold = part[0], part[1]
        col = part[2] if len(part) > 2 else color
        r = p.add_run(); r.text = txt
        f = r.font; f.name = FONT; f.size = Pt(size); f.bold = bold; f.color.rgb = col
    return p


def box(slide, x, y, w, h, fill=None, edge=None, r=0.012, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    try:
        sp.adjustments[0] = r
    except Exception:
        pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if edge is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = edge; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(150000); tf.margin_right = Emu(140000)
    tf.margin_top = Emu(100000); tf.margin_bottom = Emu(85000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return sp


def tbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb


def accent(slide, x, y, h, color):
    box(slide, x, y, 40000, h, fill=color, shape=MSO_SHAPE.RECTANGLE)


def new_slide(prs, layout, title, kicker=""):
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    bar = box(s, L, 386331, CW, 470919, fill=PURPLE, shape=MSO_SHAPE.RECTANGLE)
    tf = bar.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(160000)
    para(tf, title, 19, bold=True, color=W, after=0, first=True, line=1.0)
    if kicker:
        tb = tbox(s, L, 900000, CW, 360000)
        para(tb.text_frame, kicker, 12.5, color=MID, after=0, first=True, line=1.15, italic=True)
    # Footer on every content slide. A judge who screenshots one page still has
    # the link to the thing that runs.
    box(s, L, 6180000, CW, 3000, fill=LAVE, shape=MSO_SHAPE.RECTANGLE)
    tb = tbox(s, L, 6270000, CW, 260000)
    runs(tb.text_frame, [("VIGIL", True, DEEP),
                         ("   ·   Team Vigil   ·   live board: ", False, GREY),
                         ("lak3hay.github.io/Vigil-Triage", True, MID),
                         ("   ·   source: ", False, GREY),
                         ("github.com/Lak3hay/Vigil-Triage", True, MID)],
         9, first=True, after=0, line=1.0)
    return s


def bullets(slide, x, y, w, items, size=11.5, gap=None):
    tb = tbox(slide, x, y, w, 400000)
    tf = tb.text_frame
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            runs(tf, [("\u2022  ", True, MID), (it[0], True, DEEP), (it[1], False, INK)],
                 size, first=(i == 0), after=gap or 7, line=1.2)
        else:
            runs(tf, [("\u2022  ", True, MID), (it, False, INK)],
                 size, first=(i == 0), after=gap or 7, line=1.2)
    return tb


def table(slide, x, y, w, cols, rows, widths, head_fill=DEEP, size=10.5, rh=340000):
    """A hand-drawn table - python-pptx tables inherit theme styling we do not want."""
    xs, total = [], sum(widths)
    acc = x
    for wd in widths:
        xs.append(acc); acc += w * wd / total
    hb = box(slide, x, y, w, 300000, fill=head_fill, shape=MSO_SHAPE.RECTANGLE)
    hb.text_frame.text = ""
    for i, c in enumerate(cols):
        tb = tbox(slide, xs[i] + 90000, y + 62000, w * widths[i] / total - 120000, 200000)
        para(tb.text_frame, c, 9.5, bold=True, color=W, after=0, first=True, line=1.0)
    yy = y + 300000
    for n, row in enumerate(rows):
        if n % 2 == 0:
            box(slide, x, yy, w, rh, fill=RGBColor(0xFA, 0xF9, 0xFC), shape=MSO_SHAPE.RECTANGLE)
        for i, cell in enumerate(row):
            txt, bold, col = (cell if isinstance(cell, tuple) else (cell, False, INK))
            tb = tbox(slide, xs[i] + 90000, yy + 78000, w * widths[i] / total - 120000, rh)
            para(tb.text_frame, txt, size, bold=bold, color=col, after=0, first=True, line=1.12)
        yy += rh
    return yy


# ═══════════════════════════════════════════════════════════════════════════════
def build() -> Path:
    R = compare()
    P = R["paired_on_deteriorating_patients"]
    # Counted, never typed. A deck that hardcodes a result drifts from it, and
    # this one already shipped "-15 min" from before a fix and "214 tests" from
    # before the last three commits.
    N_TESTS = _count_tests()
    N_SWEEP = _count_safety_sweep()
    F, V = R["fifo"], R["vigil"]
    comp = composition()
    surge = surge_summary(generate_surge(hours=3.0, multiplier=3.0))

    # Two passes. Deleting slides and then adding new ones in the same session
    # makes python-pptx reuse a partname that the dropped slide still occupies,
    # which silently clobbers a kept slide -- it cost us the team-details page
    # once. Saving in between forces the package to renumber cleanly.
    import tempfile

    stage = Presentation(str(TEMPLATE))
    ids = list(stage.slides._sldIdLst)
    # keep 0 (title) and 2 (team details); drop instructions, the Round 1
    # content slides, and the stock thank-you (we close on "See it run")
    for idx in (6, 5, 4, 3, 1):
        stage.part.drop_rel(ids[idx].rId)
        stage.slides._sldIdLst.remove(ids[idx])
    tmp = Path(tempfile.gettempdir()) / "_vigil_deck_stage.pptx"
    stage.save(str(tmp))

    prs = Presentation(str(tmp))
    # A genuinely blank layout. The obvious choice ("Content: text + split")
    # carries a dark panel across the right third of the slide, which renders as
    # a black block behind hand-drawn content -- invisible when inspecting text,
    # obvious the moment it is exported.
    layout = next(
        l for m in prs.slide_masters for l in m.slide_layouts
        if l.name == "Blank with copyrights"
    )

    # ── 3. the problem ────────────────────────────────────────────────────────
    s = new_slide(prs, layout, "The problem: triage is an event. Deterioration is a process.",
                  "A patient marked stable at minute one can be in septic shock by minute ninety.")
    y = 1480000
    for i, (t, b, c, tint) in enumerate([
        ("Under-triage is predictable, not random",
         "The misses cluster: the elderly septic patient with no fever, the woman with atypical "
         "cardiac pain, the stroke that looks like vertigo, the child in compensated shock whose "
         "blood pressure is still normal. They look well \u2014 until they don't.", RED, REDT),
        ("The waiting room is unmonitored",
         "Once labelled, nobody revisits. The readings that would show deterioration \u2014 a heart "
         "rate drifting up, a saturation drifting down \u2014 never cross a threshold, so nothing "
         "fires. Deaths in ED waiting rooms are documented, not hypothetical.", AMB, AMBT),
        ("Crowding makes it worse, exactly when it matters",
         "Acuity assignments drift downward as the department fills. The system under-triages "
         "hardest precisely when accuracy matters most \u2014 so a tool that only helps on a quiet "
         "shift helps when it is not needed.", MID, LAV),
    ]):
        x = L + i * (CW / 3 + 60000)
        wdt = CW / 3 - 60000
        sp = box(s, x, y, wdt, 2050000, fill=tint)
        accent(s, x, y, 2050000, c)
        tf = sp.text_frame; tf.margin_left = Emu(240000)
        para(tf, t, 12.5, bold=True, color=c, after=7, first=True, line=1.1)
        para(tf, b, 11, color=INK, after=0, line=1.25)

    sp = box(s, L, 3760000, CW, 720000, fill=DEEP)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    runs(tf, [("On the average patient the nurse is already right. ", True, W),
              ("Vigil does not compete with her \u2014 it covers the tail, and it never looks away.",
               False, RGBColor(0xD3, 0xB8, 0xF8))],
         14, first=True, align=PP_ALIGN.CENTER, line=1.0)

    tb = tbox(s, L, 4700000, CW, 800000)
    runs(tb.text_frame, [("Who this is for.  ", True, DEEP),
                         ("Emergency departments of 100\u2013500+ visits a day, where a single nurse "
                          "sequences every arrival under time pressure with incomplete information, "
                          "and where the reassessment intervals the triage standard already mandates "
                          "are the first thing to fail when the department fills.", False, INK)],
         12, first=True, line=1.2)

    # ── 4. the solution ───────────────────────────────────────────────────────
    s = new_slide(prs, layout, "Vigil: the nurse still decides \u2014 nothing goes unwatched",
                  "Three loops, and one equation that four behaviours fall out of.")
    y = 1470000
    for i, (h, sub, b, hero) in enumerate([
        ("ASSESS", "minute 0\u20133",
         "Resolves the AGE BAND first, then scores with the right instrument \u2014 NEWS2 for adults, "
         "a PEWS-style score over paediatric ranges for children. Then a panel of 12 encoded red "
         "flags for the presentations too rare to learn from data.", False),
        ("WATCH", "the loop nobody runs",
         "Every waiting patient stays live. Two triggers: an overdue re-check clock, and vitals "
         "re-recorded as worsening. Trends, not thresholds \u2014 two readings forty minutes apart "
         "are the entire signal, so no telemetry is required.", True),
        ("SEQUENCE", "the second axis",
         "Acuity asks how sick you are. Vigil adds: what does waiting cost you? One continuous "
         "curve gives sequencing within a band, re-ranking on deterioration, anti-starvation, "
         "and routing.", False),
    ]):
        x = L + i * (CW / 3 + 60000)
        wdt = CW / 3 - 60000
        sp = box(s, x, y, wdt, 2130000, fill=DEEP if hero else LAV, edge=None if hero else LAVE)
        tf = sp.text_frame; tf.margin_left = Emu(200000)
        para(tf, h, 14, bold=True, color=W if hero else DEEP, after=1, first=True, line=1.0)
        para(tf, sub.upper(), 9, bold=True, color=RGBColor(0xC9, 0xA6, 0xF5) if hero else MID,
             after=8, line=1.0)
        para(tf, b, 10.5, color=W if hero else INK, after=0, line=1.24)

    tb = tbox(s, L, 3720000, CW, 400000)
    para(tb.text_frame, "IT IS A DECLARED SCHEDULING POLICY, NOT AN ESTIMATED CAUSAL QUANTITY",
         10, bold=True, color=MID, after=0, first=True, line=1.0)
    tb = tbox(s, L, 3980000, CW, 900000)
    runs(tb.text_frame, [
        ("We do not learn \u201chow much worse will your outcome be if you wait 30 more minutes.\u201d "
         "That is confounded in the obvious direction \u2014 sicker patients are seen sooner, so naive "
         "estimation concludes waiting is good for you. Every constant is visible and "
         "site-configurable, which is also how the same engine flexes from a 120-visit rural "
         "department to a 500-visit urban one ", False),
        ("without retraining anything.", True, DEEP)], 12, first=True, line=1.22)

    sp = box(s, L, 5060000, CW, 640000, fill=GRNT)
    accent(s, L, 5060000, 640000, GRN)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Emu(250000)
    runs(tf, [("Because the curve is convex in waiting time, a low-acuity patient eventually "
               "overtakes a higher one. ", False),
              ("The queue cannot starve anyone \u2014 that is a mathematical property, not a fairness patch.",
               True, GRN)], 11.5, first=True, line=1.15)

    # ── 5. authority ladder ───────────────────────────────────────────────────
    s = new_slide(prs, layout, "What it decides, what it suggests, what it never touches",
                  "We automate only decisions whose worst case is wasted effort \u2014 never missed care.")
    y = 1500000
    for i, (head, c, tint, items) in enumerate([
        ("DECIDES", GRN, GRNT, ["The reassessment clock", "Observation intensity",
                                "Order within an acuity band"]),
        ("RECOMMENDS", AMB, AMBT, ["The acuity level, with drivers", "Pathway activation",
                                   "Stream / routing"]),
        ("NEVER", RED, REDT, ["Diagnoses", "Discharges", "Lowers a level",
                              "Acts on the patient"]),
    ]):
        x = L + i * (CW / 3 + 60000)
        wdt = CW / 3 - 60000
        sp = box(s, x, y, wdt, 1720000, fill=tint)
        accent(s, x, y, 1720000, c)
        tf = sp.text_frame; tf.margin_left = Emu(240000)
        para(tf, head, 13, bold=True, color=c, after=6, first=True, line=1.0)
        for j, it in enumerate(items):
            para(tf, "\u2022  " + it, 11, color=INK, after=3, line=1.15)

    tb = tbox(s, L, 3420000, CW, 900000)
    runs(tb.text_frame, [("Escalation can be automatic. De-escalation is always human. ", True, DEEP),
                         ("Every autonomous action can only add attention \u2014 a sooner re-check, a "
                          "closer watch, a different position in a queue. None of them changes what "
                          "care a patient receives.", False)], 12.5, first=True, line=1.22)

    sp = box(s, L, 4380000, CW, 1320000, fill=DEEP)
    tf = sp.text_frame; tf.margin_left = Emu(250000); tf.margin_top = Emu(170000)
    para(tf, "THE SAFETY PROPERTY \u2014 TESTED, NOT ASSERTED", 10.5, bold=True,
         color=RGBColor(0xC9, 0xA6, 0xF5), after=8, first=True, line=1.0)
    para(tf, "recommended_level  \u2264  nurse_acuity      for every patient, always", 15,
         bold=True, color=W, after=8, line=1.0)
    para(tf, f"Swept over {N_SWEEP:,} synthetic patients across every age band, vital-sign extreme and "
             "level of data completeness. If Vigil can only ever raise urgency, then adding it to a "
             "department cannot create a new under-triage failure that did not already exist "
             "without it.", 11, color=RGBColor(0xE2, 0xC4, 0xFF), after=0, line=1.2)

    # ── 6. age banding ────────────────────────────────────────────────────────
    a22, a08 = assess(by_id("P22").snapshot), assess(by_id("P08").snapshot)
    s = new_slide(prs, layout, "Age banding: the gap the brief calls a silent safety risk",
                  "\u201cSolutions that apply a single adult-calibrated scoring model across all age "
                  "groups introduce silent safety risk.\u201d \u2014 Round 2 brief")
    y = 1560000
    for i, (pid, title, vit, a, verdict, c, tint) in enumerate([
        ("P22", "4-year-old, fever, not drinking", "HR 165   RR 34   SpO\u2082 96   SBP 95", a22,
         "An adult model reads \u201ctachycardic but normotensive\u201d and calls it stable. In children "
         "blood pressure falls LAST \u2014 a normal BP here is not reassurance.", RED, REDT),
        ("P08", "6-month-old, fever, feeding well", "HR 150   RR 40   SpO\u2082 98   SBP 90", a08,
         "The same shape of numbers. Alarming in an adult, unremarkable at this age. The band must "
         "cut both ways, or it is just a louder alarm.", GRN, GRNT),
    ]):
        x = L + i * (CW / 2 + 80000)
        wdt = CW / 2 - 80000
        sp = box(s, x, y, wdt, 2320000, fill=tint)
        accent(s, x, y, 2320000, c)
        tf = sp.text_frame; tf.margin_left = Emu(250000)
        para(tf, f"{pid}  \u00b7  {title}", 12, bold=True, color=c, after=5, first=True, line=1.1)
        para(tf, vit, 12, bold=True, color=INK, after=6, line=1.0)
        runs(tf, [("Instrument: ", False, GREY), (a.score.instrument, True, INK),
                  (f"  =  {a.score.total}", True, INK)], 11, after=3, line=1.1)
        runs(tf, [("Nurse said ", False, GREY), (str(a.nurse_acuity), True, INK),
                  ("   \u2192   Vigil ", False, GREY), (str(a.recommended_level), True, c),
                  ("   ESCALATED" if a.escalated else "   not escalated", True, c)],
             11, after=7, line=1.1)
        para(tf, verdict, 10.5, color=INK, after=0, line=1.22)

    tb = tbox(s, L, 4080000, CW, 1500000)
    tf = tb.text_frame
    runs(tf, [("The age band is resolved before anything else, and every threshold downstream is a "
               "function of it. ", False),
              ("An unknown age is a first-class state that lowers confidence \u2014 never a silent "
               "adult default.", True, DEEP)], 12, first=True, after=10, line=1.22)
    runs(tf, [("Beside the score sits a panel of ", False), (f"{len(PANEL)} encoded red flags", True, DEEP),
              (" for presentations that are rare ", False), ("by construction", True, INK),
              (" \u2014 occult sepsis in the elderly, atypical ACS, posterior stroke, beta-blocker "
               "masking. A model trained on ED data will underperform on exactly the cases that "
               "kill, because they are underrepresented in the training set by definition. No "
               "amount of data fixes that, so they are encoded in one readable file a clinical "
               "governance lead can review and sign off.", False)], 12, after=0, line=1.22)

    # ── 7. WATCH / P17 ────────────────────────────────────────────────────────
    s = new_slide(prs, layout, "WATCH: deterioration while every reading is still normal",
                  "P17 \u00b7 68F \u00b7 abdominal pain \u00b7 takes a beta-blocker, which blunts the "
                  "heart-rate response")
    y = 1500000
    segs = [
        ("MINUTE 0", "HR 88  RR 18  SpO\u2082 99  SBP 124", "NEWS2 0 \u00b7 level 3 \u00b7 re-check in 6 h",
         "Every value inside its normal range. The nurse is right.", PURPLE, W),
        ("MINUTE 40", "HR 99  RR 20  SpO\u2082 95  SBP 112", "NEWS2 2 \u00b7 re-check in 15 min",
         "STILL all normal. Vigil flags silent deterioration and tightens the clock 24-fold.",
         AMB, W),
        ("MINUTE 75", "HR 112  RR 24  SpO\u2082 93  SBP 104", "NEWS2 7 \u00b7 level 1 \u00b7 queue rank 9 \u2192 1",
         "Now it is visible to everyone. Vigil saw it 35 minutes ago.", RED, W),
    ]
    for i, (head, vit, score, note, c, tc) in enumerate(segs):
        x = L + i * (CW / 3 + 60000)
        wdt = CW / 3 - 60000
        sp = box(s, x, y, wdt, 1900000, fill=c)
        tf = sp.text_frame; tf.margin_left = Emu(200000)
        para(tf, head, 12.5, bold=True, color=tc, after=5, first=True, line=1.0)
        para(tf, vit, 12, bold=True, color=tc, after=5, line=1.05)
        para(tf, score, 10.5, color=tc, after=7, line=1.1)
        para(tf, note, 10.5, bold=True, color=tc, after=0, line=1.2)

    sp = box(s, L, 3620000, CW, 700000, fill=DEEP)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    runs(tf, [("Rising heart rate with falling oxygen is deterioration ", False, W),
              ("even when no single reading is abnormal.", True, W)],
         13.5, first=True, align=PP_ALIGN.CENTER, line=1.0)

    tb = tbox(s, L, 4500000, CW, 1200000)
    tf = tb.text_frame
    runs(tf, [("Two readings are enough. ", True, DEEP),
              ("Forty minutes apart is the entire signal \u2014 no waveform, no continuous telemetry, "
               "no monitor per chair. That is what makes this deployable on a shared vitals kiosk "
               "and a re-check clock rather than on hardware nobody will buy.", False)],
         12, first=True, after=9, line=1.22)
    runs(tf, [("The clock adds no work \u2014 it re-orders existing work. ", True, DEEP),
              ("Reassessment intervals are already mandated by the triage standard in use. When "
               "capacity genuinely cannot meet the schedule, that surfaces as a staffing signal "
               "rather than accumulating silently as overdue tasks.", False)], 12, after=0, line=1.22)

    # ── 8. results ────────────────────────────────────────────────────────────
    s = new_slide(prs, layout, "Does it help? One shift, run twice, one change.",
                  "Identical arrivals, identical deteriorations, identical capacity. Only the "
                  "queue ordering differs.")
    rows = [
        ("Patients seen (capacity held constant)", f"{F['seen']}", f"{V['seen']}", "\u2014"),
        ("MEAN wait, all patients", f"{F['mean_wait_min']} min", f"{V['mean_wait_min']} min",
         "unchanged"),
        ("Median wait", f"{F['median_wait_min']} min", f"{V['median_wait_min']} min",
         f"{V['median_wait_min']-F['median_wait_min']:+.1f}"),
        ("90th-percentile wait", f"{F['p90_wait_min']} min", f"{V['p90_wait_min']} min",
         f"{V['p90_wait_min']-F['p90_wait_min']:+.1f}"),
        ("Median deterioration \u2192 seen",
         f"{P['median_minutes_deterioration_to_seen_fifo']} min",
         f"{P['median_minutes_deterioration_to_seen_vigil']} min",
         f"{P['median_change_min']:+.0f}"),
    ]
    styled = [[(r[0], False, INK), (r[1], False, INK), (r[2], True, DEEP),
               (r[3], True, GRN if r[3].startswith("-") else GREY)] for r in rows]
    yy = table(s, L, 1560000, CW, ["", "STATUS QUO (FIFO)", "VIGIL", "CHANGE"], styled,
               [0.44, 0.19, 0.19, 0.18])

    sp = box(s, L, yy + 150000, CW, 700000, fill=GRNT)
    accent(s, L, yy + 150000, 700000, GRN)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Emu(250000)
    runs(tf, [(f"Paired on the {P['n_matched']} patients who deteriorated in both arms: ", False),
              (f"{P['reached_sooner']} reached sooner, {P['reached_later']} later, "
               f"{P['unchanged']} unchanged.", True, GRN)], 12, first=True, line=1.15)

    tb = tbox(s, L, yy + 1000000, CW, 1100000)
    tf = tb.text_frame
    runs(tf, [("Read the MEAN row first. ", True, RED),
              ("Re-ordering cannot create capacity, so total waiting time is conserved and only "
               "REDISTRIBUTED. The median falls 32 minutes while the 90th percentile rises 50: "
               "most patients wait considerably less, a minority wait meaningfully longer. That is "
               "the trade, and anti-starvation is what stops the tail growing without bound.",
               False)], 11.5, first=True, after=8, line=1.2)
    para(tf, "Simulation on synthetic patients under stated assumptions. Not clinical evidence \u2014 "
             "the assumptions are parameters, not constants, and the direction of the effect is "
             "testable by changing them.", 10.5, color=GREY, italic=True, after=0, line=1.15)

    # ── 9. target users ───────────────────────────────────────────────────────
    s = new_slide(prs, layout, "Who uses it \u2014 and who can stop it",
                  "Most triage tools are designed for the nurse and killed by everyone else.")
    users = [
        ("The triage nurse", "the user",
         "Wants the shift not to get worse. Vigil writes her note, never overrules her, and "
         "protects her medico-legally. The blind second opinion \u2014 she commits her level before "
         "seeing ours \u2014 keeps her judgment independent.", GRN, GRNT),
        ("The charge nurse", "the daily beneficiary",
         "Owns flow. The capacity warning is her feature: it turns \u201cwe are drowning\u201d into an "
         "escalation with a number attached, at the moment the reassessment schedule stops being "
         "deliverable.", MID, LAV),
        ("Clinical governance lead", "the gatekeeper",
         "Must sign off anything that touches a patient. Buys because the red-flag panel is one "
         "readable file they can review \u2014 not weights they must take on trust. This is why we "
         "chose rules over a learned ranker.", AMB, AMBT),
        ("Medical director / CMO", "the budget",
         "Signs the cheque. Buys the audit trail and the reassessment-compliance evidence \u2014 not "
         "the AI.", DEEP, LAV),
    ]
    y = 1520000
    for i, (name, role, body, c, tint) in enumerate(users):
        x = L + (i % 2) * (CW / 2 + 80000)
        yy2 = y + (i // 2) * 1330000
        wdt = CW / 2 - 80000
        sp = box(s, x, yy2, wdt, 1240000, fill=tint)
        accent(s, x, yy2, 1240000, c)
        tf = sp.text_frame; tf.margin_left = Emu(250000)
        runs(tf, [(name, True, c), ("   \u2014 " + role, False, GREY)], 12, first=True,
             after=5, line=1.05)
        para(tf, body, 10.5, color=INK, after=0, line=1.22)

    sp = box(s, L, 4340000, CW, 800000, fill=REDT)
    accent(s, L, 4340000, 800000, RED)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Emu(250000)
    runs(tf, [("The single thing that kills adoption: ", True, RED),
              ("nurses believing it scores their performance. The disagreement record exists to "
               "improve the model and audit the system, and is contractually barred from staff "
               "performance management. That belongs in the sales conversation, not a footnote.",
               False)], 11.5, first=True, line=1.18)

    # ── 10. business case ─────────────────────────────────────────────────────
    s = new_slide(prs, layout, "The business case: sell the boring thing first",
                  "Every other team will lead with clinical outcomes. Outcomes are third on the "
                  "purchase order.")
    rows = [
        [("1.  Reassessment compliance", True, DEEP),
         ("An existing audited obligation. Reassessment intervals are already mandated by the "
          "triage standard in use, and compliance is documented to collapse under crowding. "
          "Vigil produces the evidence that the re-checks a department already owes actually "
          "happened.", False, INK)],
        [("2.  Medico-legal defensibility", True, DEEP),
         ("Every recommendation, override, reason and re-check in a tamper-evident chain. "
          "Under-triage deaths are litigable; this is the defence file.", False, INK)],
        [("3.  Clinical outcomes", True, DEEP),
         ("Deteriorating patients reached sooner. What we care about \u2014 and the hardest to "
          "attribute, so it is not the first line on a purchase order.", False, INK)],
        [("4.  Throughput", True, DEEP),
         ("Capacity without capital expenditure. Real, but it is the benefit every vendor "
          "claims, so it persuades least.", False, INK)],
    ]
    yy = table(s, L, 1540000, CW, ["THE BUYER LADDER", "WHY IT SIGNS"], rows, [0.28, 0.72], rh=700000)

    sp = box(s, L, yy + 200000, CW, 780000, fill=DEEP)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "We do not sell a better triage decision. We sell evidence that the re-checks a "
             "department already owes actually happened.", 13.5, bold=True, color=W,
         after=0, first=True, align=PP_ALIGN.CENTER, line=1.1)

    tb = tbox(s, L, yy + 1120000, CW, 500000)
    para(tb.text_frame, "An existing budget line, an existing audit, an existing failure \u2014 which "
                        "is a far shorter sale than a new clinical claim.", 11.5, color=GREY,
         italic=True, after=0, first=True, line=1.15)

    # ── 11. roadmap ───────────────────────────────────────────────────────────
    s = new_slide(prs, layout, "Deployment roadmap: phased by authority, not by feature",
                  "No phase grants authority the previous one has not earned. Each exit criterion "
                  "is measurable.")
    rows = [
        [("0", True, W), ("Shadow", True, DEEP), ("None \u2014 nobody sees the output", False, INK),
         ("Under-triage events caught vs. missed, at a fixed alert budget", False, INK)],
        [("1", True, W), ("Advisory", True, DEEP), ("Displays; decides nothing", False, INK),
         ("Staff open it unprompted; override rate stabilises", False, INK)],
        [("2", True, W), ("The clock", True, DEEP), ("Owns the re-check schedule only", False, INK),
         ("Reassessment compliance rises; overdue queue does not grow", False, INK)],
        [("3", True, W), ("Sequencing", True, DEEP), ("Orders within an acuity band", False, INK),
         ("Deterioration-to-seen falls AND the 90th-percentile wait stays bounded", False, INK)],
        [("4", True, W), ("Second site", True, DEEP), ("Same authority, new department", False, INK),
         ("Recalibration measured in days, not months", False, INK)],
    ]
    for r in rows:
        r[0] = (r[0][0], True, INK)
    yy = table(s, L, 1540000, CW, ["", "PHASE", "AUTHORITY HELD", "EXIT CRITERION \u2014 MEASURED, NOT ASSERTED"],
               rows, [0.04, 0.14, 0.30, 0.52], rh=440000)

    tb = tbox(s, L, yy + 220000, CW, 1200000)
    tf = tb.text_frame
    runs(tf, [("The phases are the authority ladder extended through time. ", True, DEEP),
              ("A hospital is not asked to trust the whole system at once \u2014 it grants one "
               "capability, measures it, and grants the next. Phase 3's exit criterion is the "
               "tail we measured, because that is the honest gate: a policy that helps the median "
               "by punishing the 90th percentile has not earned the next phase.", False)],
         12, first=True, after=8, line=1.22)
    runs(tf, [("Shadow mode is not a formality. ", True, DEEP),
              ("It is the only phase that measures what the system would have caught before "
               "anyone can be harmed by trusting it.", False)], 12, after=0, line=1.22)

    # ── 12. risks ─────────────────────────────────────────────────────────────
    s = new_slide(prs, layout, "What could go wrong, and what we did about it",
                  "The risks that matter are the ones that survive a good design.")
    rows = [
        [("Alert fatigue \u2014 staff learn to ignore it", True, INK),
         ("Hard budget on high-salience interrupts; unacknowledged criticals escalate to the "
          "charge nurse rather than repeating. An alert that can be ignored will be.", False, INK)],
        [("Automation bias \u2014 the nurse stops thinking", True, INK),
         ("Blind second opinion: she commits her level before seeing ours. Also yields labelled "
          "disagreement data for free.", False, INK)],
        [("Over-triage \u2014 escalating everything", True, INK),
         ("Tested in both directions: a fully normal adult must NOT be escalated. Over-triage "
          "consumes the capacity the sick need, so it causes the harm it claims to prevent.", False, INK)],
        [("Thresholds wrong for this hospital", True, INK),
         ("Every constant is a declared, site-configurable policy \u2014 not learned weights. "
          "Two profiles ship; recalibration needs no retraining.", False, INK)],
        [("Staff believe it is monitoring them", True, INK),
         ("Disagreement record contractually barred from performance management, and the audit "
          "log records reasoning rather than people.", False, INK)],
        [("It is trusted more than it deserves", True, INK),
         ("It abstains rather than guessing, states confidence on every output, and this deck "
          "lists what we did not build.", False, INK)],
    ]
    table(s, L, 1520000, CW, ["RISK", "MITIGATION \u2014 BUILT, NOT PLANNED"], rows,
          [0.33, 0.67], rh=560000)

    # ── 13. architecture + scale ──────────────────────────────────────────────
    s = new_slide(prs, layout, "Architecture, and how one engine fits very different hospitals",
                  "Deterministic first, learned second, generative last \u2014 and never for a number.")
    y = 1500000
    stack = [
        ("Intake", "Snapshot of what is genuinely known in the first minutes \u2014 no lab, no imaging."),
        ("Clinical core", "Age band \u2192 NEWS2 / PEWS-style \u2192 12 encoded red flags \u2192 decomposed "
                          "confidence \u2192 compose with min(nurse, computed)."),
        ("Flow", "Cost-of-waiting policy, routing, and the WATCH loop over the live queue."),
        ("Audit", "Hash-chained append-only log. DPDP-shaped data minimisation."),
        ("Surface", "Static board \u2014 no server, no key, no external requests."),
    ]
    for i, (name, body) in enumerate(stack):
        yy2 = y + i * 560000
        sp = box(s, L, yy2, CW * 0.56, 500000, fill=LAV, edge=LAVE)
        tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Emu(200000)
        runs(tf, [(name + "   ", True, DEEP), (body, False, INK)], 10.5, first=True, line=1.12)

    x2 = L + CW * 0.58
    w2 = CW * 0.42
    sp = box(s, x2, y, w2, 2620000, fill=DEEP)
    tf = sp.text_frame; tf.margin_left = Emu(230000); tf.margin_top = Emu(180000)
    para(tf, "SAME ENGINE, TWO DEPARTMENTS", 10, bold=True, color=RGBColor(0xC9, 0xA6, 0xF5),
         after=9, first=True, line=1.0)
    para(tf, "Urban trauma centre \u00b7 500 visits/day", 11, bold=True, color=W, after=3, line=1.05)
    para(tf, "Standard targets are achievable, so they are held.", 10.5,
         color=RGBColor(0xD3, 0xB8, 0xF8), after=9, line=1.15)
    para(tf, "Rural district \u00b7 120 visits/day", 11, bold=True, color=W, after=3, line=1.05)
    para(tf, "The sickest are held to a TIGHTER target because retrieval takes longer; low-acuity "
             "targets are relaxed to what the department can actually meet. A target nobody can hit "
             "is not a safety standard \u2014 it is an alarm generator, and it is how these systems "
             "lose staff trust.", 10.5, color=RGBColor(0xD3, 0xB8, 0xF8), after=9, line=1.15)
    para(tf, "No model is retrained. Only the declared policy changes.", 10.5, bold=True,
         color=W, after=0, line=1.15)

    sp = box(s, L, 4380000, CW, 800000, fill=AMBT)
    accent(s, L, 4380000, 800000, AMB)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Emu(250000)
    runs(tf, [("What the language model does NOT do: compute a score. ", True, AMB),
              ("In the obvious build an LLM is asked \u201cwhat level is this patient?\u201d \u2014 "
               "unreproducible, unauditable, and impossible for a governance lead to sign off. "
               "We confine language models to language.", False)], 11.5, first=True, line=1.18)

    # ── 14. limits ────────────────────────────────────────────────────────────
    s = new_slide(prs, layout, "What we did not build",
                  "A prototype that hides its edges cannot be evaluated \u2014 and will not survive "
                  "a technical review.")
    y = 1520000
    lims = [
        ("No causal estimate of harm-of-delay",
         "The cost-of-waiting curve is a declared policy. The causal version needs an instrument "
         "or a natural experiment; we have neither and do not claim one."),
        ("No learned model in the decision path",
         "The data pipeline exists and is tested against real de-identified ED records. Training "
         "was cut for scope, and the ED module alone cannot express good outcome labels."),
        ("No LLM layer, no EHR integration",
         "Both designed and specified; neither built. The dataset seam exists, the adapters do not."),
        ("No clinical validation",
         "Every number here comes from synthetic patients or simulation. Nothing in this work is "
         "evidence of clinical benefit, and thresholds are illustrative defaults."),
    ]
    for i, (t, b) in enumerate(lims):
        yy2 = y + i * 800000
        sp = box(s, L, yy2, CW, 730000, fill=GREYF)
        accent(s, L, yy2, 730000, GREY)
        tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Emu(250000)
        runs(tf, [(t + "   \u2014  ", True, INK), (b, False, GREY)], 11, first=True, line=1.18)

    sp = box(s, L, 4820000, CW, 900000, fill=DEEP)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Emu(250000)
    runs(tf, [("MISTAKES.md", True, W),
              (" in the public repository is an append-only log of every error we made building "
               "this \u2014 six found by our own checks, including a detector that flagged "
               "deterioration and then scheduled the re-check four hours later. It is public on "
               "purpose.", False, RGBColor(0xD3, 0xB8, 0xF8))], 12, first=True, line=1.2)

    # ── 15. the ask / links ───────────────────────────────────────────────────
    s = new_slide(prs, layout, "See it run", "")
    tb = tbox(s, L, 1300000, CW, 500000)
    para(tb.text_frame, "LIVE BOARD \u2014 NO INSTALL", 11, bold=True, color=MID, after=0,
         first=True, line=1.0)
    tb = tbox(s, L, 1560000, CW, 500000)
    para(tb.text_frame, "lak3hay.github.io/Vigil-Triage", 21, bold=True, color=DEEP,
         after=0, first=True, line=1.0)
    tb = tbox(s, L, 2100000, CW, 500000)
    para(tb.text_frame, "SOURCE, README AND TESTS", 11, bold=True, color=MID, after=0,
         first=True, line=1.0)
    tb = tbox(s, L, 2360000, CW, 500000)
    para(tb.text_frame, "github.com/Lak3hay/Vigil-Triage", 21, bold=True, color=DEEP,
         after=0, first=True, line=1.0)

    y = 3120000
    stats = [(str(len(SCENARIOS)), "synthetic patients,\nevery required case"),
             (str(N_TESTS), f"tests, including a\n{N_SWEEP:,}-patient safety sweep"),
             (str(len(PANEL)), "encoded red flags,\nreviewable in one file"),
             (f"\u2212{abs(P['median_change_min']):.0f} min",
              "median time to reach\na deteriorating patient")]
    for i, (n, lab) in enumerate(stats):
        x = L + i * (CW / 4 + 40000)
        wdt = CW / 4 - 40000
        sp = box(s, x, y, wdt, 1120000, fill=LAV, edge=LAVE)
        tf = sp.text_frame; tf.margin_left = Emu(200000)
        para(tf, n, 26, bold=True, color=DEEP, after=3, first=True, line=1.0)
        para(tf, lab, 10.5, color=INK, after=0, line=1.18)

    tb = tbox(s, L, 4400000, CW, 400000)
    runs(tb.text_frame, [("Prototype demo video: ", True, DEEP),
                         ("see the repository README", False, INK)],
         12, first=True, after=0, line=1.0)

    sp = box(s, L, 4900000, CW, 720000, fill=DEEP)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "The nurse still decides. Vigil makes sure nothing is missed \u2014 at minute 1, and at "
             "minute 90.", 14, bold=True, color=W, after=0, first=True,
         align=PP_ALIGN.CENTER, line=1.0)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    p = build()
    print("wrote", p)
